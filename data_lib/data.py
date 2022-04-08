# library imports
import os
import codecs
import json
import win32com.client
import datetime as dt
import pandas as pd
import xlwings as xw
import fiscalyear as fy
from .utils import *
from sqlalchemy import create_engine
from pptx import Presentation


class DB:
    def __init__(self, db_config):
        self._cfg = db_config

    def set_engine(self):
        """
        Create a connection to the database based on the database type
        :return: The engine object.
        """
        if self._cfg["type"] == "postgres":
            return create_engine(
                self._cfg["connstring"].format(
                    self._cfg["user"], self._cfg["password"], self._cfg["server"], self._cfg["database"]
                )
            )
        elif self._cfg["type"] == "mssql":
            return create_engine(
                self._cfg["connstring"].format(self._cfg["server"], self._cfg["database"])
            )

        elif self._cfg["type"] == "oracle":
            return create_engine(
                self._cfg["connstring"].format(
                    self._cfg["user"], self._cfg["password"], self._cfg["server"], self._cfg["database"]
                )
            )
        elif self._cfg["type"] == "mysql":
            return create_engine(
                self._cfg["connstring"].format(
                    self._cfg["user"], self._cfg["password"], self._cfg["server"], self._cfg["database"]
                )
            )


class Data:
    def __init__(self, config):
        """
        The __init__ function is called when an instance of the class is created.
        :param config: The config object that was created in the previous step
        """
        self._cfg = config
        self._output_path = None
        self._load_path = None
        self._query_path = None
        self._input_path = None
        self._pass_path = None
        self._today = dt.datetime.today() - pd.DateOffset(months=1)
        self._monthname_long = self._today.strftime("%B")
        self._monthname_short = self._today.strftime("%b")
        self._monthyear = self._today.strftime("%m%y")
        self._fiscal_year = None
        self._current_quarter = None
        self._fiscal_quarter = None
        self._set_fiscal_year()

    def set_path(self, set_custom_path=False, path_type=None, custom_path=None):
        """
        :param set_custom_path: If you want to set a custom path, set this to True, defaults to False
        (optional)
        :param path_type: The type of path you want to set
        :param custom_path: the path to the folder where the files are located
        """

        if not set_custom_path:
            self._input_path = self._cfg.path["input"]
            self._load_path = self._cfg.path["load"]
            self._output_path = self._cfg.path["output"]
            self._query_path = self._cfg.path["queries"]
            self._pass_path = self._cfg.path["tracker"]
        else:
            if path_type == "output":
                self._output_path = custom_path
            elif path_type == "load":
                self._load_path = custom_path
            elif path_type == "query":
                self._query_path = custom_path
            elif path_type == "input":
                self._input_path = custom_path

    def _set_fiscal_year(self, f_year=None, s_month=11, s_day=1, s_year="previous"):
        """
        :param f_year: The year you want to set as the fiscal year
        :param s_month: The starting month of the fiscal year, defaults to 11 (optional)
        :param s_day: The day of the month that the fiscal year starts, defaults to 1 (optional)
        :param s_year: The start year of the fiscal year, defaults to previous (optional)
        """

        if f_year is None:
            if dt.datetime.today().month < 11:
                f_year = dt.datetime.today().year - 1
            else:
                f_year = dt.datetime.today().year

        self._fiscal_year = fy.FiscalYear(f_year)
        fy.setup_fiscal_calendar(
            start_year=s_year, start_month=s_month, start_day=s_day
        )
        self._fiscal_quarter = str(
            fy.FiscalQuarter.current().prev_fiscal_quarter
        ).split(" ")[0]
        self._current_quarter = str(
            fy.FiscalQuarter.current().prev_fiscal_quarter
        ).split(" ")[1]

    def external_query(self, query_name, db_engine):
        """
        :param query_name: The name of the query you want to run
        :return: A dataframe
        """

        query_path = self._query_path + "{0}.sql".format(query_name)
        with open(query_path) as get:
            query = get.read()
        return pd.read_sql_query(query, db_engine)

    def internal_query(self, query_name, db_engine):
        """
        :param query_name: The name of the query to run
        :return: A dataframe
        """

        return pd.read_sql_query(query_name, db_engine)

    def execute_query(self, query_name, db_engine):
        """
        Execute a query using the engine

        :param query_name: The name of the query to execute
        """
        with db_engine as connection:
            connection.execute(query_name)

    def file_query(self, file_name, engine_reader=None, read_sheet=0):
        """
        :param file_name: The name of the file to be read
        :param engine_reader: The engine to use for reading the Excel file
        :param read_sheet: The sheet number to read, defaults to 0 (optional)
        :return: A dataframe
        """

        file = self._input_path + file_name
        if engine_reader is None:
            return pd.read_excel(file, sheet_name=read_sheet, index_col=0)
        else:
            return pd.read_excel(
                file, sheet_name=read_sheet, index_col=0, engine=engine_reader
            )

    def _password_tracker(self, request_date, output_file_name, password):
        """
        :param request_date: The date and time when the request was made
        :param output_file_name: The name of the file that will be generated
        :param password: The password to protect the output file. If None, the file will be unprotected
        """

        file_name = "password_tracker.json"
        list_obj = []

        if os.path.isfile(self._pass_path + file_name) is False:
            with open(self._pass_path + file_name, "w"):
                pass

        # Read JSON file
        if os.stat(self._pass_path + file_name).st_size != 0:
            with open(self._pass_path + file_name) as get_content:
                list_obj = json.load(get_content)

        list_obj.append(
            {
                "Request Date": str(request_date),
                "File Name": output_file_name,
                "Path": self._output_path,
                "Password": "Unprotected" if password is None else password,
            }
        )

        with open(self._pass_path + file_name, "w") as json_file:
            json.dump(list_obj, json_file, indent=4, separators=(",", ": "))

    def _password_share(self, password, email_subject):
        """
        :param password: The password to be shared
        :param email_subject: The subject of the email that will be sent to the user
        """

        outlook = win32com.client.Dispatch("Outlook.Application").GetNamespace("MAPI")
        sent_items = outlook.GetDefaultFolder(5)
        messages = sent_items.Items
        message = messages.GetLast()

        print(password, email_subject, message)

        pass

    def _draft_email(self, recipient, file_name, password, folder_search):
        """
        :param recipient: The email address of the recipient
        :param file_name: The name of the file that will be attached to the email
        :param password: The password to protect the file with
        :param folder_search: The name of the folder to search for the email. If None, the search will
        be done in the inbox
        :return: None
        """

        outlook = win32com.client.Dispatch("Outlook.Application").GetNamespace("MAPI")
        if folder_search is not None:
            inbox = outlook.GetDefaultFolder(6).Folders.Item(
                folder_search
            )  # inbox in sub-folders
        else:
            inbox = outlook.GetDefaultFolder(6)  # inbox emails

        # instantiate outlook to send email
        const = win32com.client.constants
        ol_mail_item = 0x0
        obj = win32com.client.Dispatch("Outlook.Application")

        # get received emails
        messages = inbox.Items
        message = messages.GetLast()
        sender_name = message.Sender.GetExchangeUser().name.split(",")[-1].strip()
        sender_address = message.Sender.GetExchangeUser().PrimarySmtpAddress
        sender_cc_address = ";".join(
            [
                item.AddressEntry.GetExchangeUser().PrimarySmtpAddress
                for item in message.Recipients
                if item.AddressEntry.GetExchangeUser().PrimarySmtpAddress
                != "andres.garcia.fernandez@hp.com"
            ]
        )

        # set email signature
        signature_htm = os.path.join(
            os.environ["USERPROFILE"], "AppData\\Roaming\\Microsoft\\Signatures\\HP.htm"
        )
        html_file = codecs.open(signature_htm, "r", "utf-8", errors="ignore")
        email_signature = html_file.read()
        html_file.close()

        # html email body and subject

        # get original email subject
        email_subject = message.Subject

        # create new email items
        new_mail = obj.CreateItem(ol_mail_item)
        new_mail.Subject = email_subject
        new_mail.BodyFormat = 2

        if password is not None:
            email_body = """
                    <HTML>
                        <BODY style="font-family:HP Simplified Light;font-size:14.5px;">
                            <p>Hi {0},</p>
                            <p>Please find attached the report requested.</p>
                            <p>
                            The attached file is password protected, password will be shared on the following email.
                            </p>
                            <p><em style="color:red">*** Delete Before Sending Email ***</em>
                            <strong>Password:</strong> {1} Copy password to send it on a separate email.
                            <em style="color:red">*** Delete Before Sending Email ***</em></p>
                            <p>Kind Regards,</p> 
                        </BODY>
                    </HTML>""".format(
                sender_name, password
            )
        else:
            email_body = """
                <HTML>
                    <BODY style="font-family:HP Simplified Light;font-size:14.5px;">
                        <p>Hi {0},</p>
                        <p>Please find attached the report requested.</p>
                        <p>Kind Regards,</p> 
                    </BODY>
                </HTML>""".format(
                sender_name
            )

        # validate that email will be sent to correct recipient
        if sender_address == recipient:  # based on the subject replying to email
            reply_all = message.ReplyAll()
            new_mail.HTMLBody = email_body + email_signature + reply_all.HTMLBody
            new_mail.To = sender_address
            new_mail.CC = sender_cc_address
            attachment = self._output_path + file_name + ".xlsx"
            new_mail.Attachments.Add(Source=attachment)
            new_mail.save()

        return None

    def _file_saver(
        self, data, file_name, protect_file, draft_email, requester, email_folder
    ):
        """
        :param data: The data to be written to the Excel file
        :param file_name: The name of the file to be saved
        :param protect_file: If True, the file will be password protected
        :param draft_email: If True, the email will be saved as a draft. If False, the email will be
        sent immediately
        :param requester: The email address of the requester
        :param email_folder: The folder where the email will be saved
        """

        # default password
        file_password = None

        # remove existing files before save
        ut.file_cleaner("{0}{1}.xlsx".format(self._output_path, file_name))

        # workbook / sheet variables
        app = xw.App(visible=False)
        wb = xw.Book()
        ws = wb.sheets[0]

        # excel data header formatting
        ws.range("A1").options(pd.DataFrame, index=False).value = ut.number_to_string(
            data, "ID"
        )
        header_format = ws.range("A1").expand("right")
        header_format.color = self._cfg.file_format["header_bg_color"]
        header_format.api.Font.Name = self._cfg.file_format["font_name"]
        header_format.api.Font.Color = self._cfg.file_format["header_font_color"]
        header_format.api.Font.Bold = True
        header_format.api.Font.Size = self._cfg.file_format["header_font_size"]

        # excel data content formatting
        data_format = ws.range("A2").current_region
        data_format.api.Font.Name = self._cfg.file_format["font_name"]
        data_format.api.Font.Size = self._cfg.file_format["content_font_size"]

        # save password protect file if needed
        if protect_file:
            file_password = "HPI" + str(self._today.strftime("%I%M%S"))
            wb.api.SaveAs(self._output_path + file_name, Password=file_password)

        else:
            wb.api.SaveAs(self._output_path + file_name)

        app.quit()

        # update password tracker
        self._password_tracker(
            self._today,
            file_name,
            file_password,
        )

        # prompt
        if draft_email:
            self._draft_email(requester, file_name, file_password, email_folder)

    def export_data(
        self,
        odf,
        custom_filename=None,
        protect_file=False,
        draft_email=False,
        requester=None,
        email_folder=None,
    ):
        """
        :param odf: the dataframe to be exported
        :param custom_filename: If you want to save the file with a custom name, you can do so
        :param protect_file: If True, the file will be protected with a password, defaults to False
        (optional)
        :param draft_email: If True, the email will be saved as a draft in Outlook, defaults to False
        (optional)
        :param requester: The email address of the person who requested the report
        :param email_folder: The name of the folder in which to save the email. If None, the email will
        be saved in the root folder
        """

        # assign file name
        if custom_filename is None:
            file_name = "Output Report " + self._today.strftime("%Y-%m-%d")
        else:
            file_name = custom_filename

        # save file
        self._file_saver(
            odf, file_name, protect_file, draft_email, requester, email_folder
        )

    def ppt_export(self, file_type, template_type, org=None):
        """
        :param file_type: "HPI" or "L1 ORG"
        :param template_type: The template type is either "HPI" or "L1 ORG"
        :param org: The name of the organization you want to export
        """

        # determine output file HPI Total or L1 Org Dashboard
        if file_type == "HPI":
            output_file = "HPI DEI Dashboard {0} {1}".format(
                self._current_quarter, self._fiscal_quarter
            )
        elif file_type == "L1 ORG":
            output_file = "DEI {0} Dashboard {1}".format(org, self._monthyear)

        # set slides to be updated
        prs = Presentation(
            "../Templates/HP_Presentation_Template_" + template_type + ".pptx"
        )
        slide_1 = prs.slides[0]
        slide_4 = prs.slides[3]
        slide_6 = prs.slides[5]
        slide_7 = prs.slides[6]
        slide_8 = prs.slides[7]
        slide_9 = prs.slides[8]

        # set placeholders to be updated
        prs_sub_1 = slide_1.placeholders[1]
        prs_sub_4 = slide_4.placeholders[0]
        prs_sub_6 = slide_6.placeholders[0]
        prs_sub_7 = slide_7.placeholders[0]
        prs_sub_8 = slide_8.placeholders[0]
        prs_sub_9 = slide_9.placeholders[0]

        # update placeholders
        if file_type == "HPI":
            prs_sub_1.text = "As of {0} month, end/{1}".format(
                self._monthname_long, self._current_quarter
            )
        elif file_type == "L1 ORG":
            prs_sub_1.text = "{0} (As of {1} month, end/{2})".format(
                org, self._monthname_long, self._current_quarter
            )

        prs_sub_4.text = (
            self._fiscal_quarter + " " + self._current_quarter + " Headcount"
        )

        prs_sub_6.text = "{0}/{1} Status to Diversity Targets (Company Level)".format(
            self._monthname_short, self._current_quarter
        )
        prs_sub_7.text = "{0}/{1} Active Headcount by Organization / MRU".format(
            self._monthname_short, self._current_quarter
        )
        prs_sub_8.text = (
            "{0}/{1} Active Headcount by Organization / MRU (Absolute values)".format(
                self._monthname_short, self._current_quarter
            )
        )
        prs_sub_9.text = "{0}/{1} US Ethnic Groups Not Self-Identified HC by Organization / MRU".format(
            self._monthname_short, self._current_quarter
        )

        # save updated template
        prs.save(
            self.__output_path
            + "Quarterly Dashboards\\"
            + file_type
            + "\\"
            + output_file
            + ".pptx"
        )
