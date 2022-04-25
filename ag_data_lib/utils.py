import os
import datetime as dt
from pptx import Presentation
from traitlets.config import Config
import nbformat as nbf
from nbconvert.exporters import HTMLExporter
from nbconvert.preprocessors import TagRemovePreprocessor

# delete existing file
def file_cleaner(file_name):
    """
    Deletes any old file before storing an updated file

    Parameters
    ----------
    file_name : str
        The file location and name to be deleted
    """

    if os.path.exists(file_name):
        os.remove(file_name)


# create default folders
def create_path():
    """
    It creates a list of paths to be created, then it creates them
    :return: A list of paths to the directories created.
    """
    output_path = []
    root = os.path.dirname(os.getcwd())
    target_dir = "02_data"
    paths = ["01_input_files", "02_output_files", "03_query_files", "04_password_files"]
    
    for path in paths:
        if not os.path.isdir(root + f"\\{target_dir}\\{path}"):   
            os.mkdir(root + f"\\{target_dir}\\{path}")         
            output_path.append(root + f"\\{target_dir}\\{path}")
            
    return output_path


# trims whitespaces
def column_trim(df):
    """
    Trims any trailing spaces on string dataframe columns

    Parameters
    ----------
    df : dataframe
        The dataframe that needs column trims

    Returns
    -------
    dataframe
        a dataframe with all column strings trimmed
    """

    trim_strings = lambda x: x.strip() if isinstance(x, str) else x
    return df.applymap(trim_strings)


def number_to_string(df, column_name):
    """
    Converts a column of numbers to strings

    Parameters
    ----------
    df : pandas.DataFrame
        The dataframe to be converted
    column_name : str
        The column name to be converted

    Returns
    -------
    pandas.DataFrame
        The converted dataframe
    """
    if column_name in df.columns:
        df[column_name] = "'" + df[column_name]
        return df
    else:
        return df


def password_generator(co_key):
    """
    It takes a string as an argument and returns a string that is the concatenation of the argument and
    the current time
    
    :param co_key: This is the company key that is used to identify the company
    :return: the concatenation of the co_key and the current time.
    """
    return co_key + str(dt.datetime.today().strftime("%I%M%S"))


def notebook_to_html(notebook_path, html_path):
    """
    It takes a Jupyter notebook and converts it to an HTML file
    
    :param notebook_path: The path to the notebook you want to convert
    :param html_path: The path to the output HTML file
    """
    # Setup config
    cfg = Config()

    cfg.TemplateExporter.exclude_input = True
    cfg.TemplateExporter.exclude_input_prompt = True
    cfg.TagRemovePreprocessor.enabled = True

    # Configure and run out exporter
    cfg.HTMLExporter.preprocessors = ["nbconvert.preprocessors.TagRemovePreprocessor"]

    exporter = HTMLExporter(config=cfg)
    exporter.register_preprocessor(TagRemovePreprocessor(config=cfg), True)

    # Configure and run our exporter - returns a tuple - first element with html,
    # second with notebook metadata
    try:
        output = HTMLExporter(config=cfg).from_filename(notebook_path)
        html_output_name = notebook_path.rsplit('.', 1)[0] + '.html'
        # Write to output html file
        with open(html_path + f"/{html_output_name}", "w", encoding="utf8") as f:
            f.write(output[0])
    except FileNotFoundError:
        print("Notebook not found. Review the Notebook path.")


def ppt_identifier(input_file, slide_number):
    """
    Identifies a ppt single slide elements

    Parameters
    ----------
    input_file : str
        The ppt template directory location and name
    slide_number : int
        The slide number
    """

    prs = Presentation(input)
    slide = prs.slides[slide_number - 1]
    for shape in slide.shapes:
        print(
            "id: %s, name: %s, , type: %s"
            % (shape.shape_id, shape.name, shape.shape_type)
        )


def create_folder_tree(root_folder, sub_dir_folder, folder_list, path):

    # Create directory
    try:
        # Create target Directory
        os.mkdir(path + root_folder)
        os.mkdir(path + root_folder + "/" + sub_dir_folder)
        print(f"Directory {root_folder} and {sub_dir_folder} created")
    except FileExistsError:
        os.mkdir(path + root_folder + "/" + sub_dir_folder)
        print("Directory ", root_folder, " already exists")

    for folder in folder_list:
        try:
            # Create target Directory
            os.mkdir(path + root_folder + "/" + sub_dir_folder + "/" + folder)
            print(f"Directory {folder} created")
        except FileExistsError:
            print("Directory ", folder, " already exists")


def ppt_analyzer(input_path, output_path):
    """
    Review a PowerPoint presentation to indicate each of ppt elements

    Parameters
    ----------
    input_path : str
        The ppt template directory location and name
    output_path : str
        The file output to avoid overwrite the original template
    """

    prs = Presentation(input_path)

    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = "Title for Layout {}".format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))

        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if "Title" not in shape.text:
                        shape.text = "Placeholder index:{} type:{}".format(
                            phf.idx, shape.name
                        )
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                print("{} {}".format(phf.idx, shape.name))
    prs.save(output_path)
